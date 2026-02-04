#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
協業プロジェクト計画 PowerPoint追加スライド生成スクリプト（改善版）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_collaboration_slides():
    """既存のプレゼンテーションに協業計画スライドを追加"""

    # 既存のプレゼンテーションを読み込み
    prs = Presentation("物流ソリューション提案書_ヤマエ久野.pptx")

    # スライド追加
    add_divider_slide(prs, "協業プロジェクト計画")
    add_collaboration_approach(prs)
    add_10_themes_overview(prs)
    add_themes_detail_1(prs)
    add_themes_detail_2(prs)
    add_themes_detail_3(prs)
    add_project_timeline(prs)
    add_project_structure(prs)
    add_standard_process(prs)
    add_cumulative_effects(prs)
    add_success_factors(prs)
    add_next_steps_detail(prs)

    return prs


def add_divider_slide(prs, title_text):
    """セクション区切りスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色的な効果
    left = Inches(0)
    top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(26, 84, 144)
    shape.line.fill.background()

    # タイトル
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    return slide


def add_collaboration_approach(prs):
    """協業アプローチ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    # タイトル
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "協業アプローチ：テーマ別段階的実現"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)

    # リード文（背景色付き）
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9)
    height = Inches(0.7)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
    shape.line.color.rgb = RGBColor(26, 84, 144)
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = "現場実務に深く入り込み、テーマ別に段階的に成果を創出。\nクイックウィンで早期効果を実証し、確実に物流改革を実現します。"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.alignment = PP_ALIGN.CENTER

    # 基本方針（左側ボックス）
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(4.3)
    height = Inches(4.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 248, 252)
    shape.line.color.rgb = RGBColor(26, 84, 144)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = "■ 基本方針"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(10)

    approaches = [
        "✓ テーマ別プロジェクトで",
        "  段階的に成果創出",
        "",
        "✓ 現場実務に深く入り込み、",
        "  実態を理解",
        "",
        "✓ クイックウィンで",
        "  早期効果を実証",
        "",
        "✓ パイロット→検証→横展開",
        "  の確実な進め方",
        "",
        "✓ 貴社メンバーと協働、",
        "  ノウハウ移転",
    ]

    for item in approaches:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("✓"):
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(13)
        p.space_before = Pt(2)

    # プロジェクト期間（右側ボックス）
    left = Inches(5.2)
    top = Inches(2.0)
    width = Inches(4.3)
    height = Inches(4.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(252, 245, 245)
    shape.line.color.rgb = RGBColor(192, 57, 43)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = "■ プロジェクト期間：36ヶ月"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    # Year 1
    p = tf.add_paragraph()
    p.text = "Year 1（Month 1-6）"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "現状可視化、クイックウィン"
    p.font.size = Pt(13)
    p.space_before = Pt(2)

    p = tf.add_paragraph()
    p.text = "効果：150-300百万円/年"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(2)

    # Year 2
    p = tf.add_paragraph()
    p.text = "Year 2（Month 7-18）"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "基盤構築、パイロット実施"
    p.font.size = Pt(13)
    p.space_before = Pt(2)

    p = tf.add_paragraph()
    p.text = "効果：500-700百万円/年"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(2)

    # Year 3
    p = tf.add_paragraph()
    p.text = "Year 3（Month 19-36）"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "全社展開、自走体制確立"
    p.font.size = Pt(13)
    p.space_before = Pt(2)

    p = tf.add_paragraph()
    p.text = "効果：870-1,240百万円/年"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(2)

    return slide


def add_10_themes_overview(prs):
    """10テーマ概要"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "10の重点テーマ"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    themes = [
        "1. 在庫可視化・実態把握（2-3ヶ月）",
        "   → 滞留在庫特定、即時削減200-500百万円",
        "",
        "2. 需要予測精度向上（3-4ヶ月）",
        "   → 予測誤差±15%→±5-8%、在庫削減300-600百万円",
        "",
        "3. 適正在庫基準策定（2-3ヶ月）",
        "   → 在庫削減800-1,500百万円、回転日数7-10日短縮",
        "",
        "4. 配送ルート最適化（3-4ヶ月）",
        "   → 配送コスト削減500-800百万円/年",
        "",
        "5. 倉庫オペレーション改善（3-4ヶ月/拠点）",
        "   → 生産性20-30%向上、人件費50-100百万円削減",
    ]

    for item in themes:
        p = tf.add_paragraph()
        p.text = item
        if item and item[0].isdigit():
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide


def add_themes_detail_1(prs):
    """テーマ詳細1"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "10の重点テーマ（続き）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    themes = [
        "6. 物流コスト可視化（2-3ヶ月）",
        "   → コスト構造把握、改善ターゲット特定",
        "",
        "7. VMI導入パイロット（6-9ヶ月）",
        "   → 在庫20-30%削減、欠品率50%削減",
        "",
        "8. IoT・デジタル化パイロット（4-6ヶ月）",
        "   → リアルタイム可視化、配送効率化",
        "",
        "9. 物流人材育成（継続）",
        "   → スキル向上、多能工化、属人化解消",
        "",
        "10. KPIダッシュボード構築（3-4ヶ月）",
        "    → 経営可視化、データドリブン経営",
    ]

    for item in themes:
        p = tf.add_paragraph()
        p.text = item
        if item and item[0].isdigit():
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(6)

    return slide


def add_themes_detail_2(prs):
    """テーマ詳細2（実施順序）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "テーマ実施タイミング"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Phase 1：現状分析・クイックウィン（Month 1-6）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(8)

    phase1 = [
        "• テーマ1：在庫可視化",
        "• テーマ6：物流コスト可視化",
        "• テーマ5：倉庫改善（1拠点）",
        "効果：150-300百万円/年",
    ]

    for item in phase1:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("効果"):
            p.font.bold = True
        p.font.size = Pt(15)
        p.space_before = Pt(4)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Phase 2：基盤構築・パイロット（Month 6-18）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(8)
    p.space_before = Pt(12)

    phase2 = [
        "• テーマ2：需要予測、テーマ3：適正在庫",
        "• テーマ4：配送最適化",
        "• テーマ7：VMI、テーマ8：IoT",
        "• テーマ10：KPIダッシュボード",
        "効果：500-700百万円/年（累計）",
    ]

    for item in phase2:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("効果"):
            p.font.bold = True
        p.font.size = Pt(15)
        p.space_before = Pt(4)

    p = tf.add_paragraph()
    p.text = ""

    p = tf.add_paragraph()
    p.text = "Phase 3：全社展開（Month 18-36）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(8)
    p.space_before = Pt(12)

    phase3 = [
        "• システム導入（WMS, TMS）",
        "• 全拠点展開、定着化",
        "効果：870-1,240百万円/年（フル効果）",
    ]

    for item in phase3:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("効果"):
            p.font.bold = True
        p.font.size = Pt(15)
        p.space_before = Pt(4)

    return slide


def add_themes_detail_3(prs):
    """標準活動プロセス（8ステップ）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "標準活動プロセス（8ステップ）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    steps = [
        "Step 1：キックオフ・計画策定（Week 1-2）",
        "Step 2：現状調査・データ収集（Week 3-6）",
        "Step 3：分析・課題整理（Week 7-10）",
        "Step 4：改善案設計（Week 11-14）",
        "Step 5：承認・意思決定（Week 15）",
        "Step 6：パイロット実施（Week 16-24）",
        "Step 7：効果検証・横展開判断（Week 25-26）",
        "Step 8：横展開・定着化（Week 27-52）",
    ]

    for i, step in enumerate(steps):
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)
        if i in [4, 6]:  # 意思決定ポイントを強調
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "各テーマでこの標準プロセスを実施"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_before = Pt(20)

    return slide


def add_project_timeline(prs):
    """プロジェクトタイムライン"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "プロジェクト全体タイムライン（36ヶ月）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    timeline = [
        "Month 0-1：準備・キックオフ",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "体制構築、詳細計画、データ収集環境整備",
        "",
        "Month 1-6：現状分析・クイックウィン",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "在庫可視化、倉庫改善、早期成果創出",
        "効果：150-300百万円/年",
        "",
        "Month 6-18：基盤構築・パイロット",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "需要予測、配送最適化、VMI/IoTパイロット",
        "システム選定・要件定義",
        "効果：500-700百万円/年（累計）",
        "",
        "Month 18-36：全社展開・定着化",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "システム導入、全拠点展開、自走体制確立",
        "効果：870-1,240百万円/年（フル効果）",
    ]

    for item in timeline:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("Month"):
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        elif item.startswith("効果"):
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
        elif item.startswith("━"):
            p.font.size = Pt(14)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide


def add_project_structure(prs):
    """プロジェクト推進体制"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "プロジェクト推進体制"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    structure = [
        "ステアリングコミッティ（月1回）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "役割：重要意思決定、進捗確認",
        "メンバー：社長、役員、PMO",
        "",
        "プロジェクトマネジメントオフィス（週1回）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "役割：全体統括、進捗管理、課題管理",
        "メンバー：プロジェクトリーダー、各テーマリーダー",
        "",
        "テーマ別ワーキンググループ（週1-2回）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "役割：テーマ別の詳細検討・実行",
        "構成：10チーム",
        "メンバー：現場責任者・実務担当者（3-5名）",
        "        コンサルタント（1-2名）",
    ]

    for item in structure:
        p = tf.add_paragraph()
        p.text = item
        if item.endswith("回）"):
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
        elif item.startswith("━"):
            p.font.size = Pt(14)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide


def add_standard_process(prs):
    """標準プロセス図式"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "活動プロセスの流れ"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    process_flow = [
        "現状調査・データ収集",
        "  ↓",
        "分析・課題整理（根本原因分析）",
        "  ↓",
        "改善案設計（To-Be、効果試算）",
        "  ↓",
        "【意思決定】承認・予算確保",
        "  ↓",
        "パイロット実施（小規模で検証）",
        "  ↓",
        "効果検証・改善調整",
        "  ↓",
        "【Go/No Go判断】横展開判断",
        "  ↓",
        "全社展開・定着化",
        "  ↓",
        "PDCAサイクル・継続改善",
    ]

    for item in process_flow:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("【"):
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        elif item == "  ↓":
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.size = Pt(16)
        p.space_before = Pt(2)

    return slide


def add_cumulative_effects(prs):
    """累積効果の推移"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "累積効果の推移"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    effects = [
        "Year 1（Month 1-12）クイックウィン創出期",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• 在庫削減：200-500百万円",
        "• 倉庫改善：30-50百万円/年",
        "• 配送初期改善：50-100百万円/年",
        "累計効果：150-300百万円/年（目標の15-25%）",
        "Year 1 CF：△50-100百万円",
        "",
        "Year 2（Month 13-24）本格展開開始期",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• 需要予測・適正在庫：+400-700百万円",
        "• 配送最適化：+200-300百万円/年",
        "• VMI/倉庫展開：+150-250百万円/年",
        "累計効果：500-700百万円/年（目標の50-60%）",
        "Year 2 CF：+50-200百万円（回収開始）",
        "",
        "Year 3（Month 25-36）フル効果達成期",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• システム稼働：+200-300百万円/年",
        "• 全拠点展開：+100-150百万円/年",
        "• VMI/IoT拡大：+70-150百万円/年",
        "累計効果：870-1,240百万円/年（目標100%）",
        "Year 3 CF：+720-1,090百万円",
    ]

    for item in effects:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("Year"):
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        elif item.startswith("累計効果") or "CF：" in item:
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
        elif item.startswith("━"):
            p.font.size = Pt(14)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(3)

    return slide


def add_success_factors(prs):
    """成功の鍵"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "成功の5つの鍵"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    factors = [
        "1. 経営層の強いコミットメント",
        "   トップダウン推進、明確な目標、リソース確保",
        "",
        "2. 現場を巻き込んだ推進",
        "   現場の声を反映、早期成果で実感創出",
        "",
        "3. データドリブンの徹底",
        "   事実に基づく課題把握、定量効果測定",
        "",
        "4. クイックウィンの創出",
        "   早期成果実証、組織の信頼獲得",
        "",
        "5. 人材育成とノウハウ移転",
        "   貴社メンバー成長、内製化、自走体制",
    ]

    for item in factors:
        p = tf.add_paragraph()
        p.text = item
        if item and item[0].isdigit():
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(16)
        p.space_before = Pt(10)

    return slide


def add_next_steps_detail(prs):
    """次のステップ詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "次のステップ：プロジェクト開始まで"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    next_steps = [
        "Step 1：詳細ヒアリング・現地視察（2-3週間）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• 経営層・物流部門ヒアリング",
        "• 拠点視察（2-3拠点）",
        "• データ確認・収集",
        "成果物：詳細現状分析、優先テーマ案、詳細計画",
        "",
        "Step 2：提案プレゼンテーション（1週間）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• 詳細計画のご説明",
        "• 期待効果の精緻化",
        "• 体制・役割分担の確認",
        "",
        "Step 3：キックオフ（1週間）",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "• キックオフミーティング",
        "• プロジェクト体制発足",
        "• 活動開始",
    ]

    for item in next_steps:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("Step"):
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        elif item.startswith("成果物"):
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
        elif item.startswith("━"):
            p.font.size = Pt(14)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(4)

    return slide


if __name__ == "__main__":
    prs = add_collaboration_slides()
    prs.save("物流ソリューション提案書_ヤマエ久野_完全版.pptx")
    print("協業プロジェクト計画スライドを追加しました")
    print(f"総スライド数: {len(prs.slides)}")
