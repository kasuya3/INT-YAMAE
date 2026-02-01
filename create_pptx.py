#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ヤマエ久野株式会社 物流ソリューション提案書 PowerPoint生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """プレゼンテーションを作成"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # スライド1: タイトル
    slide = add_title_slide(prs)

    # スライド2: 目次
    slide = add_agenda_slide(prs)

    # スライド3: エグゼクティブサマリー
    slide = add_executive_summary_slide(prs)

    # スライド4-6: 現状分析（財務指標）
    slide = add_financial_analysis_1(prs)
    slide = add_financial_analysis_2(prs)
    slide = add_financial_analysis_3(prs)

    # スライド7-10: 課題分析
    slide = add_issue_1(prs)
    slide = add_issue_2(prs)
    slide = add_issue_3(prs)
    slide = add_issue_4(prs)

    # スライド11: ソリューション全体像
    slide = add_solution_overview(prs)

    # スライド12-14: 各ソリューション詳細
    slide = add_solution_1(prs)
    slide = add_solution_2(prs)
    slide = add_solution_3(prs)

    # スライド15: 投資対効果サマリー
    slide = add_roi_summary(prs)

    # スライド16: 財務改善シミュレーション
    slide = add_financial_simulation(prs)

    # スライド17-18: 収益性改善の内訳
    slide = add_profitability_improvement_1(prs)
    slide = add_profitability_improvement_2(prs)

    # スライド19: 投資計画
    slide = add_investment_plan(prs)

    # スライド20-22: 実行ロードマップ
    slide = add_roadmap_phase1(prs)
    slide = add_roadmap_phase2(prs)
    slide = add_roadmap_phase3(prs)

    # スライド23: 期待効果まとめ
    slide = add_expected_benefits(prs)

    # スライド24: 次のステップ
    slide = add_next_steps(prs)

    # スライド25: Thank you
    slide = add_thank_you_slide(prs)

    return prs


def add_title_slide(prs):
    """タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    # タイトル
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "物流システムソリューション提案書"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)

    # サブタイトル
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "株式会社ヤマエ久野 御中"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(52, 73, 94)

    # 日付
    left = Inches(1)
    top = Inches(6)
    width = Inches(8)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "2026年2月"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(127, 140, 141)

    return slide


def add_agenda_slide(prs):
    """目次スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "目次"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    items = [
        "1. エグゼクティブサマリー",
        "2. 現状分析：財務指標から見た経営課題",
        "3. 根本原因：物流課題の詳細分析",
        "4. 提案ソリューション全体像",
        "5. 各ソリューション詳細",
        "6. 投資対効果・財務改善シミュレーション",
        "7. 実行ロードマップ",
        "8. 期待効果まとめ"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(12)

    return slide


def add_executive_summary_slide(prs):
    """エグゼクティブサマリー"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "エグゼクティブサマリー"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # Key message
    p = tf.add_paragraph()
    p.text = "経営課題と提案"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    items = [
        "キャッシュフロー悪化：営業CFマージン -3.12%",
        "収益性の低迷：経常利益率 1.73%、総利益率 7.38%",
        "財務健全性の問題：流動比率 91.06%",
        "根本原因：在庫14,000百万円、物流コスト推定35,000百万円",
        "",
        "提案：3つの物流システムソリューション",
        "投資額：970-1,260百万円",
        "年間効果：870-1,240百万円",
        "投資回収：約1.1-1.4年",
        "経常利益率改善：1.73% → 3.0-3.5%"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.space_before = Pt(6)

    return slide


def add_financial_analysis_1(prs):
    """財務分析1：キャッシュフロー"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "現状分析1：キャッシュフロー悪化"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # TDB報告書の図表を追加（画像）
    try:
        left = Inches(0.5)
        top = Inches(1.8)
        height = Inches(4.5)
        pic = slide.shapes.add_picture('企業情報/tdb_page-43.png', left, top, height=height)
    except:
        pass

    # テキストボックスで解説
    left = Inches(6)
    top = Inches(1.8)
    width = Inches(3.5)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "重要な示唆"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)

    items = [
        "期首現金残高の減少",
        "5,811→2,767百万円",
        "",
        "営業CFの激しい変動",
        "150億円レベルの変動",
        "",
        "営業CFマージン",
        "-3.12%（マイナス）",
        "",
        "主因：棚卸資産増減",
        "△1,429百万円"
    ]

    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide


def add_financial_analysis_2(prs):
    """財務分析2：収益性指標"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "現状分析2：収益性指標の著しい低迷"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "業界標準との比較"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)

    # テーブル風にテキストで表現
    metrics = [
        ("指標", "現状", "業界標準", "評価"),
        ("売上高総利益率", "7.38%", "10-15%", "❌ 著しく低い"),
        ("売上高経常利益率", "1.73%", "3-5%", "❌ 大幅に低い"),
        ("流動比率", "91.06%", "150%以上", "❌ 低水準"),
    ]

    for row in metrics:
        p = tf.add_paragraph()
        p.text = f"{row[0]:12} {row[1]:10} {row[2]:12} {row[3]}"
        if row[0] == "指標":
            p.font.bold = True
            p.font.size = Pt(16)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "根本原因："
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(20)

    causes = [
        "調達・物流コストの高さ → 総利益率7.38%",
        "物流コスト推定7-9% → 経常利益率1.73%",
        "在庫14,000百万円 → 流動比率91.06%"
    ]

    for cause in causes:
        p = tf.add_paragraph()
        p.text = cause
        p.font.size = Pt(16)
        p.space_before = Pt(8)

    return slide


def add_financial_analysis_3(prs):
    """財務分析3：物流起因の構造"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "現状分析3：物流起因の財務悪化構造"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # 構造図（テキストで表現）
    structure = [
        "物流課題",
        "  ↓",
        "① 在庫管理の非効率",
        "  → 在庫14,000百万円（過剰・長期滞留）",
        "  → 棚卸資産増減 △1,429百万円",
        "  → 営業CFマージン -3.12%",
        "  → 流動比率 91.06%",
        "",
        "② 物流コストの増大",
        "  → 推定35,000百万円（売上比7-9%）",
        "  → 経常利益率 1.73%",
        "",
        "③ 調達・在庫ロスの発生",
        "  → FIFO管理不徹底、品質劣化",
        "  → 総利益率 7.38%",
    ]

    for line in structure:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("①") or line.startswith("②") or line.startswith("③"):
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(4)

    return slide


def add_issue_1(prs):
    """課題1：在庫管理の非効率"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "課題1：在庫管理の非効率"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "現状の問題点"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    issues = [
        "在庫水準：約14,000百万円（運転資金の圧迫）",
        "棚卸資産増減：△1,429百万円（令和7年度）",
        "需要予測精度：属人的な発注、予測システム未導入",
        "在庫回転：業界標準より7-10日長い",
        "適正在庫不明：SKU別の需要変動分析なし",
        "",
        "財務インパクト",
        "• 営業CFマージン -3.12%の主因",
        "• 運転資金14,000百万円の固定化",
        "• 在庫ロス・陳腐化による損失発生",
        "• 保管コスト・管理コストの増大"
    ]

    for item in issues:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("財務インパクト"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_before = Pt(16)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_issue_2(prs):
    """課題2：物流コストの増大"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "課題2：物流コストの増大"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "現状の問題点"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    issues = [
        "推定物流コスト：35,000百万円（売上比7-9%）",
        "※効率的企業は売上比5-6%程度",
        "",
        "非効率の要因",
        "• 配送ルート最適化なし（手作業配車）",
        "• 積載効率：推定65-70%（最適80-85%）",
        "• 配送頻度・リードタイム最適化なし",
        "• 共同配送・モーダルシフト未実施",
        "",
        "財務インパクト",
        "• 経常利益率1.73%を圧迫",
        "• 年間3,000-4,000百万円の削減余地",
        "• 固定費（倉庫・人件費）の高止まり"
    ]

    for item in issues:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("非効率の要因") or item.startswith("財務インパクト"):
            p.font.size = Pt(19)
            p.font.bold = True
            p.space_before = Pt(14)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_issue_3(prs):
    """課題3：オペレーション非効率"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "課題3：倉庫・配送オペレーション非効率"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "現状の問題点"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    issues = [
        "倉庫作業",
        "• 紙伝票による入出荷管理",
        "• ピッキングミス率：1-2%",
        "• 作業生産性：業界標準の70-80%",
        "• ロケーション管理の不徹底",
        "",
        "配送業務",
        "• ドライバー不足・残業時間の増加",
        "• 配送計画の属人化",
        "• リアルタイム進捗管理なし",
        "",
        "財務インパクト",
        "• 人件費・固定費の高止まり",
        "• 作業ミスによる追加コスト",
        "• 生産性低下による機会損失"
    ]

    for item in issues:
        p = tf.add_paragraph()
        p.text = item
        if item in ["倉庫作業", "配送業務", "財務インパクト"]:
            p.font.size = Pt(19)
            p.font.bold = True
            p.space_before = Pt(14)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_issue_4(prs):
    """課題4：収益性・財務健全性"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "課題4：収益性・財務健全性の低迷"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "物流起因の収益性問題"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    structure = [
        "総利益率 7.38%（業界10-15%）",
        "  原因1：調達物流コスト高（VMI未導入）",
        "  原因2：在庫ロス・品質劣化による値引き",
        "",
        "経常利益率 1.73%（業界3-5%）",
        "  原因：物流コスト推定35,000百万円",
        "  売上比7-9%（効率的企業は5-6%）",
        "",
        "流動比率 91.06%（健全水準150%以上）",
        "  原因：在庫14,000百万円の固定化",
        "  運転資金の圧迫",
        "",
        "改善ポテンシャル",
        "• 総利益率：7.38% → 9.0-10.0%（+1.6-2.6pt）",
        "• 経常利益率：1.73% → 3.0-3.5%（+1.3-1.8pt）",
        "• 流動比率：91.06% → 120-130%（+29-39pt）"
    ]

    for line in structure:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("総利益率") or line.startswith("経常利益率") or line.startswith("流動比率") or line.startswith("改善ポテンシャル"):
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(4)

    return slide


def add_solution_overview(prs):
    """ソリューション全体像"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "提案ソリューション全体像"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "3つの統合ソリューション"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(16)

    solutions = [
        "① 高度在庫管理システム",
        "   AI需要予測 × IMS × VMI",
        "   投資：220-280百万円 / 効果：330-470百万円/年",
        "   → CFマージン・流動比率改善",
        "",
        "② 統合物流プラットフォーム",
        "   TMS × SCM × リアルタイム可視化",
        "   投資：350-430百万円 / 効果：330-450百万円/年",
        "   → 物流コスト削減・経常利益率改善",
        "",
        "③ 物流自動化・最適化",
        "   次世代WMS × 倉庫自動化 × モーダルシフト",
        "   投資：400-550百万円 / 効果：210-320百万円/年",
        "   → 固定費削減・生産性向上",
    ]

    for item in solutions:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("①") or item.startswith("②") or item.startswith("③"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
        else:
            p.font.size = Pt(14)
        p.space_before = Pt(6)

    return slide


def add_solution_1(prs):
    """ソリューション1詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ソリューション① 高度在庫管理システム"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    details = [
        "主要機能",
        "• AI需要予測システム（機械学習・需要変動分析）",
        "• IMS統合在庫管理（リアルタイム可視化）",
        "• VMI（ベンダー管理在庫）導入",
        "• 適正在庫自動算出",
        "",
        "導入効果",
        "• 在庫削減：1,000-2,000百万円",
        "• 在庫回転日数：7-10日短縮",
        "• 欠品率：50%削減",
        "• 予測精度：±15% → ±5%",
        "",
        "投資対効果",
        "• 初期投資：220-280百万円",
        "• 年間効果：330-470百万円",
        "• ROI：0.5-0.8年",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["主要機能", "導入効果", "投資対効果"]:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
            p.space_before = Pt(16)
        else:
            p.font.size = Pt(16)
            p.space_before = Pt(6)

    return slide


def add_solution_2(prs):
    """ソリューション2詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ソリューション② 統合物流プラットフォーム"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    details = [
        "主要機能",
        "• TMS（配送ルート最適化・動的配車）",
        "• SCM可視化（サプライチェーン全体）",
        "• IoT活用（車両・貨物追跡）",
        "• KPI自動集計・分析",
        "",
        "導入効果",
        "• 配送コスト：15-20%削減",
        "• 積載効率：65-70% → 80-85%",
        "• リードタイム：20-30%短縮",
        "• 配車計画時間：80%削減",
        "",
        "投資対効果",
        "• 初期投資：350-430百万円",
        "• 年間効果：330-450百万円",
        "• ROI：0.8-1.3年",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["主要機能", "導入効果", "投資対効果"]:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
            p.space_before = Pt(16)
        else:
            p.font.size = Pt(16)
            p.space_before = Pt(6)

    return slide


def add_solution_3(prs):
    """ソリューション3詳細"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "ソリューション③ 物流自動化・最適化"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    details = [
        "主要機能・設備",
        "• 次世代WMS（ピッキング最適化）",
        "• AGV（無人搬送車）・自動仕分け機",
        "• デジタルピッキングシステム",
        "• 共同物流・モーダルシフト",
        "",
        "導入効果",
        "• ピッキング生産性：30-40%向上",
        "• 倉庫人件費：20-30%削減（80-120百万円）",
        "• 作業ミス率：1-2% → 0.2%以下",
        "• 長距離輸送コスト：10-15%削減",
        "",
        "投資対効果",
        "• 初期投資：400-550百万円",
        "• 年間効果：210-320百万円",
        "• ROI：1.3-2.6年",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["主要機能・設備", "導入効果", "投資対効果"]:
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
            p.space_before = Pt(16)
        else:
            p.font.size = Pt(16)
            p.space_before = Pt(6)

    return slide


def add_roi_summary(prs):
    """投資対効果サマリー"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "全体投資対効果サマリー"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    # テーブル風表現
    p = tf.add_paragraph()
    p.text = "ソリューション別投資計画"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(12)

    table_data = [
        ("ソリューション", "初期投資", "年間効果", "ROI"),
        ("① 高度在庫管理", "220-280百万円", "330-470百万円", "0.5-0.8年"),
        ("② 統合物流PF", "350-430百万円", "330-450百万円", "0.8-1.3年"),
        ("③ 物流自動化", "400-550百万円", "210-320百万円", "1.3-2.6年"),
        ("", "", "", ""),
        ("合計", "970-1,260百万円", "870-1,240百万円", "1.1-1.4年"),
    ]

    for row in table_data:
        if row[0] == "":
            p = tf.add_paragraph()
            p.text = ""
            continue
        p = tf.add_paragraph()
        p.text = f"{row[0]:14} {row[1]:16} {row[2]:16} {row[3]}"
        if row[0] == "ソリューション" or row[0] == "合計":
            p.font.bold = True
            p.font.size = Pt(16)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "投資回収：2年目後半には完全回収"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(20)

    return slide


def add_financial_simulation(prs):
    """財務改善シミュレーション"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "財務指標改善シミュレーション"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "システム導入後3年目（フル効果）"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)

    # テーブル
    table_data = [
        ("財務指標", "現状", "改善後", "改善幅"),
        ("営業CF", "4,653百万円", "18,000-20,000", "+13,000-15,000"),
        ("CFマージン", "-3.12%", "3.5-4.0%", "+6.5-7.0pt"),
        ("物流コスト", "35,000百万円", "31,000-32,000", "△3,000-4,000"),
        ("", "", "", ""),
        ("総利益率", "7.38%", "9.0-10.0%", "+1.6-2.6pt"),
        ("経常利益率", "1.73%", "3.0-3.5%", "+1.3-1.8pt"),
        ("流動比率", "91.06%", "120-130%", "+29-39pt"),
    ]

    for row in table_data:
        if row[0] == "":
            p = tf.add_paragraph()
            p.text = ""
            continue
        p = tf.add_paragraph()
        if row[0] == "財務指標":
            p.text = f"{row[0]:12} {row[1]:14} {row[2]:16} {row[3]}"
            p.font.bold = True
            p.font.size = Pt(15)
        else:
            p.text = f"{row[0]:12} {row[1]:14} {row[2]:16} {row[3]}"
            p.font.size = Pt(14)
        p.space_before = Pt(5)

    return slide


def add_profitability_improvement_1(prs):
    """収益性改善の内訳1"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "収益性改善の内訳①：総利益率"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "総利益率：7.38% → 9.0-10.0%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(16)

    breakdown = [
        "改善要因1：調達物流最適化",
        "  VMI・共同調達によるコスト削減",
        "  効果：+0.8-1.2ポイント",
        "",
        "改善要因2：在庫ロス削減",
        "  FIFO徹底・品質管理による値引き・廃棄削減",
        "  効果：+0.5-0.8ポイント",
        "",
        "改善要因3：物流効率化",
        "  欠品削減、配送品質向上による付加価値向上",
        "  効果：+0.3-0.6ポイント",
        "",
        "合計改善幅：+1.6-2.6ポイント",
    ]

    for line in breakdown:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("改善要因") or line.startswith("合計"):
            p.font.size = Pt(18)
            p.font.bold = True
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(6)

    return slide


def add_profitability_improvement_2(prs):
    """収益性改善の内訳2"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "収益性改善の内訳②：経常利益率・流動比率"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "経常利益率：1.73% → 3.0-3.5%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(10)

    breakdown1 = [
        "• 物流コスト削減：+0.6-0.8pt",
        "• 総利益率改善効果：+0.4-0.6pt",
        "• オペレーション効率化：+0.3-0.4pt",
        "合計：+1.3-1.8ポイント",
    ]

    for line in breakdown1:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("合計"):
            p.font.size = Pt(17)
            p.font.bold = True
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(5)

    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "流動比率：91.06% → 120-130%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(10)

    breakdown2 = [
        "• 在庫最適化：14,000 → 12,000-13,000百万円",
        "• 過剰在庫削減：1,000-2,000百万円",
        "• 運転資金サイクル改善",
        "• 営業CF改善による手元流動性向上",
        "合計：+29-39ポイント",
    ]

    for line in breakdown2:
        p = tf.add_paragraph()
        p.text = line
        if line.startswith("合計"):
            p.font.size = Pt(17)
            p.font.bold = True
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(5)

    return slide


def add_investment_plan(prs):
    """投資計画"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "段階的投資計画とCF影響"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "3年間の投資・効果計画"
    p.font.size = Pt(22)
    p.font.bold = True
    p.space_after = Pt(12)

    table_data = [
        ("期間", "投資額", "効果創出", "ネットCF"),
        ("1年目", "△500-600百万円", "+150-250百万円", "△300-400百万円"),
        ("2年目", "△400-500百万円", "+500-700百万円", "+50-200百万円"),
        ("3年目", "△100-150百万円", "+870-1,240百万円", "+720-1,090百万円"),
        ("", "", "", ""),
        ("累計", "△1,000-1,250", "+1,520-2,190", "+470-940百万円"),
    ]

    for row in table_data:
        if row[0] == "":
            p = tf.add_paragraph()
            p.text = ""
            continue
        p = tf.add_paragraph()
        p.text = f"{row[0]:8} {row[1]:18} {row[2]:18} {row[3]}"
        if row[0] == "期間" or row[0] == "累計":
            p.font.bold = True
            p.font.size = Pt(16)
        else:
            p.font.size = Pt(15)
        p.space_before = Pt(6)

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "投資回収：2年目後半に完了"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_before = Pt(20)

    p = tf.add_paragraph()
    p.text = "3年目以降：フルベネフィット創出"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_before = Pt(10)

    return slide


def add_roadmap_phase1(prs):
    """ロードマップ フェーズ1"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "実行ロードマップ：フェーズ1（0-6ヶ月）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "目標：基盤構築とクイックウィン"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(12)

    details = [
        "実施内容",
        "1. 物流実態調査（現場ヒアリング・データ分析）",
        "2. 在庫分析システム導入（AI需要予測準備）",
        "3. TMS基本機能導入開始（配送最適化）",
        "4. パイロット拠点選定（1-2拠点）",
        "",
        "投資額：150-200百万円",
        "期待効果：年間100-150百万円",
        "",
        "主要KPI",
        "• 在庫精度向上：誤差率 3% → 1%以下",
        "• 配送コスト削減：5%削減",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["実施内容", "主要KPI"]:
            p.font.size = Pt(19)
            p.font.bold = True
            p.space_before = Pt(14)
        elif item.startswith("投資額") or item.startswith("期待効果"):
            p.font.size = Pt(17)
            p.font.bold = True
            p.space_before = Pt(12)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_roadmap_phase2(prs):
    """ロードマップ フェーズ2"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "実行ロードマップ：フェーズ2（6-18ヶ月）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "目標：コアシステム導入と展開"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(12)

    details = [
        "実施内容",
        "1. AI需要予測システム本格稼働",
        "2. 統合物流プラットフォーム導入",
        "3. WMS高度化と倉庫自動化設備導入",
        "4. サプライチェーン可視化の実現",
        "",
        "投資額：600-750百万円（累計750-950百万円）",
        "期待効果：年間500-700百万円",
        "",
        "主要KPI",
        "• 在庫回転日数：7日短縮",
        "• 物流コスト率：1.5%削減",
        "• 倉庫生産性：20%向上",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["実施内容", "主要KPI"]:
            p.font.size = Pt(19)
            p.font.bold = True
            p.space_before = Pt(14)
        elif item.startswith("投資額") or item.startswith("期待効果"):
            p.font.size = Pt(17)
            p.font.bold = True
            p.space_before = Pt(12)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_roadmap_phase3(prs):
    """ロードマップ フェーズ3"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "実行ロードマップ：フェーズ3（18-36ヶ月）"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "目標：全社展開と最適化"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(12)

    details = [
        "実施内容",
        "1. 全拠点へのシステム展開完了",
        "2. 倉庫自動化の拡大展開",
        "3. 継続的改善活動（PDCA）",
        "4. 新技術導入検討（AI高度化等）",
        "",
        "投資額：200-250百万円（累計950-1,200百万円）",
        "期待効果：年間870-1,240百万円（フル効果）",
        "",
        "主要KPI",
        "• 在庫回転日数：さらに3-5日短縮",
        "• 物流コスト率：売上比6%以下",
        "• 倉庫生産性：業界トップ水準達成",
        "• 経常利益率：3.0-3.5%達成",
    ]

    for item in details:
        p = tf.add_paragraph()
        p.text = item
        if item in ["実施内容", "主要KPI"]:
            p.font.size = Pt(19)
            p.font.bold = True
            p.space_before = Pt(14)
        elif item.startswith("投資額") or item.startswith("期待効果"):
            p.font.size = Pt(17)
            p.font.bold = True
            p.space_before = Pt(12)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_expected_benefits(prs):
    """期待効果まとめ"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "期待効果まとめ"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "財務指標の改善"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(12)

    benefits = [
        "CFマージン：-3.12% → 3.5-4.0%",
        "経常利益率：1.73% → 3.0-3.5%",
        "総利益率：7.38% → 9.0-10.0%",
        "流動比率：91.06% → 120-130%",
        "",
        "業務改善効果",
        "• 在庫削減：1,000-2,000百万円",
        "• 物流コスト削減：3,000-4,000百万円/年",
        "• 運転資金解放：5-7億円",
        "• 倉庫生産性向上：30-40%",
        "• 配送効率向上：積載率80-85%",
        "",
        "経営基盤の強化",
        "• データドリブン経営の実現",
        "• サプライチェーンの可視化・最適化",
        "• 競争力強化・事業成長の基盤構築",
    ]

    for item in benefits:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("業務改善効果") or item.startswith("経営基盤"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(26, 84, 144)
            p.space_before = Pt(16)
        else:
            p.font.size = Pt(16)
            p.space_before = Pt(6)

    return slide


def add_next_steps(prs):
    """次のステップ"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "次のステップ"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "提案実現に向けて"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)
    p.space_after = Pt(16)

    steps = [
        "ステップ1：詳細ヒアリング・現地調査",
        "  • 物流拠点視察",
        "  • 現場担当者ヒアリング",
        "  • データ収集・分析",
        "  期間：2-3週間",
        "",
        "ステップ2：詳細提案書作成",
        "  • システム詳細仕様",
        "  • 投資計画詳細",
        "  • 実行計画詳細",
        "  期間：3-4週間",
        "",
        "ステップ3：プロジェクト開始",
        "  • キックオフミーティング",
        "  • プロジェクト体制構築",
        "  • フェーズ1実行開始",
    ]

    for item in steps:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("ステップ"):
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(192, 57, 43)
            p.space_before = Pt(14)
        else:
            p.font.size = Pt(15)
            p.space_before = Pt(6)

    return slide


def add_thank_you_slide(prs):
    """Thank youスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト

    # メッセージ
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "ご清聴ありがとうございました"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 84, 144)

    # サブメッセージ
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(8)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "ご質問・ご相談はお気軽にお申し付けください"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(52, 73, 94)

    return slide


if __name__ == "__main__":
    prs = create_presentation()
    prs.save("物流ソリューション提案書_ヤマエ久野.pptx")
    print("PowerPointプレゼンテーションを作成しました: 物流ソリューション提案書_ヤマエ久野.pptx")
    print(f"総スライド数: {len(prs.slides)}")
