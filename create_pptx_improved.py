#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
物流ソリューション提案書 PowerPoint生成（改善版・完全版）
見栄えを大幅に改善：レイアウト最適化、視覚要素追加、フォントサイズ調整
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE

# 色定義
COLOR_PRIMARY = RGBColor(26, 84, 144)  # ブルー
COLOR_SECONDARY = RGBColor(192, 57, 43)  # レッド
COLOR_ACCENT = RGBColor(40, 116, 166)  # ライトブルー
COLOR_BG_LIGHT = RGBColor(245, 248, 252)  # 薄いブルー背景
COLOR_BG_LIGHT2 = RGBColor(252, 245, 245)  # 薄いレッド背景
COLOR_TEXT = RGBColor(52, 73, 94)  # グレー系

def create_presentation():
    """完全版プレゼンテーション作成"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 元のスライド（1-25）
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_executive_summary_slide(prs)
    add_section_divider(prs, "現状分析")
    add_financial_analysis_compact(prs)
    add_issues_summary(prs)

    add_section_divider(prs, "提案ソリューション")
    add_solution_overview(prs)
    add_solutions_detail(prs)  # 3ソリューション を1スライドに
    add_roi_summary(prs)
    add_financial_simulation(prs)

    add_section_divider(prs, "実行計画")
    add_roadmap_overview(prs)
    add_expected_benefits(prs)

    # 協業プロジェクト計画（新規）
    add_section_divider(prs, "協業プロジェクト計画")
    add_collaboration_approach(prs)
    add_10_themes_visual(prs)
    add_project_timeline_visual(prs)
    add_standard_process_visual(prs)
    add_project_structure_visual(prs)
    add_cumulative_effects_visual(prs)
    add_success_factors_visual(prs)
    add_next_steps_visual(prs)

    add_thank_you_slide(prs)

    return prs


def add_section_divider(prs, title_text):
    """セクション区切りスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景グラデーション風（2色の矩形を重ねる）
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = COLOR_PRIMARY
    shape1.line.fill.background()

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    return slide


def add_title_slide(prs):
    """タイトルスライド（改善版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景装飾
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2), Inches(10), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 245, 250)
    shape.line.fill.background()

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.text = "物流システムソリューション提案書"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # サブタイトル
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    tf.text = "株式会社ヤマエ久野 御中"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TEXT

    # 日付
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.4))
    tf = txBox.text_frame
    tf.text = "2026年2月"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(127, 140, 141)

    return slide


def add_agenda_slide(prs):
    """目次スライド（改善版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "目次"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # コンテンツ領域
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(7)
    height = Inches(5.5)

    items = [
        "1. エグゼクティブサマリー",
        "2. 現状分析：財務指標から見た経営課題",
        "3. 提案ソリューション全体像",
        "4. 投資対効果・財務改善シミュレーション",
        "5. 実行ロードマップ",
        "6. 協業プロジェクト計画",
        "7. 期待効果まとめ",
    ]

    current_top = top
    for i, item in enumerate(items):
        # アイコン（番号の円）
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left - Inches(0.5), current_top - Inches(0.05), Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_PRIMARY
        circle.line.fill.background()

        circle_tf = circle.text_frame
        circle_p = circle_tf.paragraphs[0]
        circle_p.text = str(i + 1)
        circle_p.alignment = PP_ALIGN.CENTER
        circle_p.font.size = Pt(16)
        circle_p.font.bold = True
        circle_p.font.color.rgb = RGBColor(255, 255, 255)

        # テキスト
        txBox = slide.shapes.add_textbox(left, current_top, width, Inches(0.4))
        tf = txBox.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT

        current_top += Inches(0.7)

    return slide


def add_executive_summary_slide(prs):
    """エグゼクティブサマリー（改善版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "エグゼクティブサマリー"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 経営課題ボックス
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(4.5)
    height = Inches(2.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(252, 240, 240)
    shape.line.color.rgb = COLOR_SECONDARY
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "■ 経営課題"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(8)

    issues = [
        "• CFマージン -3.12%",
        "• 経常利益率 1.73%（低い）",
        "• 流動比率 91.06%（低い）",
        "• 在庫14,000百万円（過剰）",
        "• 物流コスト推定35,000百万円",
    ]

    for issue in issues:
        p = tf.add_paragraph()
        p.text = issue
        p.font.size = Pt(13)
        p.space_before = Pt(3)

    # 提案ソリューションボックス
    left = Inches(5.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = "■ 提案ソリューション"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    solutions = [
        "① 高度在庫管理システム",
        "② 統合物流プラットフォーム",
        "③ 物流自動化・最適化",
        "",
        "3つの統合ソリューション",
    ]

    for sol in solutions:
        p = tf.add_paragraph()
        p.text = sol
        if sol.startswith("①") or sol.startswith("②") or sol.startswith("③"):
            p.font.size = Pt(13)
            p.font.bold = True
        else:
            p.font.size = Pt(13)
        p.space_before = Pt(3)

    # 投資効果ボックス（下部全幅）
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(9)
    height = Inches(3.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 252, 245)
    shape.line.color.rgb = RGBColor(39, 174, 96)
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = "■ 投資対効果"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)
    p.space_after = Pt(12)

    # 3列で表示
    roi_text = [
        "【投資額】             【年間効果】           【投資回収期間】",
        "970-1,260百万円      870-1,240百万円       約1.1-1.4年",
        "",
        "【財務改善目標】",
        "• 営業CFマージン：-3.12% → 3.5-4.0%（+6.5-7.0pt）",
        "• 経常利益率：1.73% → 3.0-3.5%（+1.3-1.8pt）",
        "• 総利益率：7.38% → 9.0-10.0%（+1.6-2.6pt）",
        "• 流動比率：91.06% → 120-130%（+29-39pt）",
    ]

    for text in roi_text:
        p = tf.add_paragraph()
        p.text = text
        if text.startswith("【"):
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_before = Pt(8)
        else:
            p.font.size = Pt(14)
            p.space_before = Pt(4)

    return slide


def add_financial_analysis_compact(prs):
    """財務分析（コンパクト版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "現状分析：財務指標から見た経営課題"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 3つのボックスを横に配置
    boxes_data = [
        {
            "title": "キャッシュフロー悪化",
            "items": [
                "現金残高",
                "5,811→2,767百万円",
                "",
                "営業CFマージン",
                "-3.12%",
                "",
                "棚卸資産増減",
                "△1,429百万円",
            ],
            "color": RGBColor(252, 240, 240),
            "line_color": COLOR_SECONDARY
        },
        {
            "title": "収益性の低迷",
            "items": [
                "総利益率",
                "7.38%（低い）",
                "",
                "経常利益率",
                "1.73%（低い）",
                "",
                "物流コスト",
                "推定7-9%",
            ],
            "color": RGBColor(255, 245, 235),
            "line_color": RGBColor(230, 126, 34)
        },
        {
            "title": "財務健全性の問題",
            "items": [
                "流動比率",
                "91.06%（低い）",
                "",
                "在庫",
                "14,000百万円",
                "",
                "回転日数",
                "業界+7-10日長い",
            ],
            "color": RGBColor(245, 240, 252),
            "line_color": RGBColor(142, 68, 173)
        }
    ]

    left_start = Inches(0.5)
    width = Inches(3)
    height = Inches(5.5)

    for i, box_data in enumerate(boxes_data):
        left = left_start + i * Inches(3.15)
        top = Inches(1.2)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_data["color"]
        shape.line.color.rgb = box_data["line_color"]
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = box_data["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = box_data["line_color"]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(15)

        for item in box_data["items"]:
            p = tf.add_paragraph()
            p.text = item
            if item and not item.startswith("•"):
                if any(c.isdigit() or c in ['-', '△', '%', '→'] for c in item):
                    p.font.size = Pt(15)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_SECONDARY
                else:
                    p.font.size = Pt(13)
                    p.font.bold = True
            else:
                p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(4)

    return slide


def add_issues_summary(prs):
    """課題サマリー"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "根本原因：4つの物流課題"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 4つの課題を2x2レイアウト
    issues = [
        {
            "num": "1",
            "title": "在庫管理の非効率",
            "detail": "在庫14,000百万円、回転日数+7-10日、予測システム未導入",
            "effect": "運転資金圧迫、CF悪化"
        },
        {
            "num": "2",
            "title": "物流コストの増大",
            "detail": "推定35,000百万円（7-9%）、配送非効率、積載率65-70%",
            "effect": "経常利益率1.73%を圧迫"
        },
        {
            "num": "3",
            "title": "オペレーション非効率",
            "detail": "紙伝票、ミス率1-2%、生産性70-80%、属人化",
            "effect": "固定費高止まり、品質問題"
        },
        {
            "num": "4",
            "title": "収益性・財務健全性",
            "detail": "総利益率7.38%、流動比率91.06%、調達コスト高",
            "effect": "競争力低下、財務リスク"
        }
    ]

    positions = [
        (Inches(0.5), Inches(1.2)),
        (Inches(5.2), Inches(1.2)),
        (Inches(0.5), Inches(4.2)),
        (Inches(5.2), Inches(4.2))
    ]

    width = Inches(4.3)
    height = Inches(2.6)

    for i, (issue, pos) in enumerate(zip(issues, positions)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_SECONDARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)

        # 番号とタイトル
        p = tf.paragraphs[0]
        p.text = f"課題{issue['num']}：{issue['title']}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(8)

        # 詳細
        p = tf.add_paragraph()
        p.text = issue['detail']
        p.font.size = Pt(12)
        p.space_before = Pt(4)

        # 影響
        p = tf.add_paragraph()
        p.text = f"→ {issue['effect']}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(8)

    return slide


def add_solution_overview(prs):
    """ソリューション全体像"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "提案ソリューション全体像"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # リード文
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "3つの統合ソリューションで、在庫・物流コスト・オペレーション効率を抜本的に改善"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # 3つのソリューション（縦に配置）
    solutions = [
        {
            "num": "①",
            "title": "高度在庫管理システム",
            "subtitle": "AI需要予測 × IMS × VMI",
            "investment": "220-280百万円",
            "effect": "330-470百万円/年",
            "roi": "0.5-0.8年",
            "target": "CFマージン・流動比率改善"
        },
        {
            "num": "②",
            "title": "統合物流プラットフォーム",
            "subtitle": "TMS × SCM × リアルタイム可視化",
            "investment": "350-430百万円",
            "effect": "330-450百万円/年",
            "roi": "0.8-1.3年",
            "target": "物流コスト削減・経常利益率改善"
        },
        {
            "num": "③",
            "title": "物流自動化・最適化",
            "subtitle": "次世代WMS × 倉庫自動化 × モーダルシフト",
            "investment": "400-550百万円",
            "effect": "210-320百万円/年",
            "roi": "1.3-2.6年",
            "target": "固定費削減・生産性向上"
        }
    ]

    top = Inches(1.8)
    height = Inches(1.6)

    for sol in solutions:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{sol['num']} {sol['title']}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = sol['subtitle']
        p.font.size = Pt(13)
        p.font.italic = True
        p.space_before = Pt(2)

        p = tf.add_paragraph()
        p.text = f"投資：{sol['investment']}  |  効果：{sol['effect']}  |  ROI：{sol['roi']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(39, 174, 96)
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"→ {sol['target']}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(4)

        top += Inches(1.75)

    return slide


def add_solutions_detail(prs):
    """3ソリューション詳細（1スライドに統合）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    tf.text = "3つのソリューション詳細"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    solutions = [
        {
            "num": "①",
            "title": "高度在庫管理",
            "functions": "AI需要予測、IMS、VMI、適正在庫算出",
            "effects": "在庫削減1,000-2,000百万円、回転日数7-10日短縮"
        },
        {
            "num": "②",
            "title": "統合物流PF",
            "functions": "TMS配送最適化、SCM可視化、IoT追跡",
            "effects": "配送コスト15-20%削減、積載率80-85%"
        },
        {
            "num": "③",
            "title": "物流自動化",
            "functions": "次世代WMS、AGV、デジタルピッキング",
            "effects": "生産性30-40%向上、人件費20-30%削減"
        }
    ]

    width = Inches(9)
    height = Inches(1.9)
    top = Inches(0.9)

    for sol in solutions:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{sol['num']} {sol['title']}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = f"機能：{sol['functions']}"
        p.font.size = Pt(13)
        p.space_before = Pt(3)

        p = tf.add_paragraph()
        p.text = f"効果：{sol['effects']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(39, 174, 96)
        p.space_before = Pt(5)

        top += Inches(2.05)

    return slide


def add_roi_summary(prs):
    """投資対効果サマリー"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "全体投資対効果サマリー"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # テーブル風の表示
    headers = ["ソリューション", "初期投資", "年間効果", "ROI"]
    rows = [
        ["① 高度在庫管理", "220-280百万円", "330-470百万円", "0.5-0.8年"],
        ["② 統合物流PF", "350-430百万円", "330-450百万円", "0.8-1.3年"],
        ["③ 物流自動化", "400-550百万円", "210-320百万円", "1.3-2.6年"],
    ]
    total = ["合計", "970-1,260百万円", "870-1,240百万円", "1.1-1.4年"]

    # ヘッダー
    left_start = Inches(0.8)
    top = Inches(1.5)
    widths = [Inches(2.5), Inches(2), Inches(2), Inches(1.5)]

    current_left = left_start
    for i, header in enumerate(headers):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            current_left, top, widths[i], Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.color.rgb = RGBColor(255, 255, 255)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        current_left += widths[i]

    # データ行
    top += Inches(0.5)
    for row in rows:
        current_left = left_start
        for i, cell in enumerate(row):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                current_left, top, widths[i], Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_BG_LIGHT
            shape.line.color.rgb = RGBColor(200, 200, 200)

            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(13)
            p.alignment = PP_ALIGN.CENTER

            current_left += widths[i]
        top += Inches(0.5)

    # 合計行
    current_left = left_start
    for i, cell in enumerate(total):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            current_left, top, widths[i], Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 250, 230)
        shape.line.color.rgb = COLOR_SECONDARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.alignment = PP_ALIGN.CENTER

        current_left += widths[i]

    # ポイント強調
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "■ 投資回収のポイント"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(39, 174, 96)
    p.space_after = Pt(10)

    points = [
        "✓ Year 2後半には投資回収完了",
        "✓ Year 3以降はフルベネフィット創出",
        "✓ 早期のクイックウィンで投資の正当性を実証",
        "✓ 段階的投資でリスクを最小化",
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(6)

    return slide


def add_financial_simulation(prs):
    """財務改善シミュレーション"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "財務指標改善シミュレーション（3年後）"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # テーブル
    headers = ["財務指標", "現状", "改善後", "改善幅"]
    data = [
        ["営業CF", "4,653百万円", "18,000-20,000", "+13,000-15,000"],
        ["CFマージン", "-3.12%", "3.5-4.0%", "+6.5-7.0pt"],
        ["", "", "", ""],  # 空行
        ["総利益率", "7.38%", "9.0-10.0%", "+1.6-2.6pt"],
        ["経常利益率", "1.73%", "3.0-3.5%", "+1.3-1.8pt"],
        ["流動比率", "91.06%", "120-130%", "+29-39pt"],
        ["", "", "", ""],  # 空行
        ["物流コスト", "35,000百万円", "31,000-32,000", "△3,000-4,000"],
    ]

    left_start = Inches(1.2)
    top = Inches(1.3)
    widths = [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8)]

    # ヘッダー
    current_left = left_start
    for i, header in enumerate(headers):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            current_left, top, widths[i], Inches(0.45)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.color.rgb = RGBColor(255, 255, 255)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        current_left += widths[i]

    # データ行
    top += Inches(0.45)
    for row in data:
        if all(cell == "" for cell in row):
            top += Inches(0.15)
            continue

        current_left = left_start
        for i, cell in enumerate(row):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                current_left, top, widths[i], Inches(0.42)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_BG_LIGHT
            shape.line.color.rgb = RGBColor(200, 200, 200)

            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(12)
            if i == 3:  # 改善幅列
                p.font.bold = True
                p.font.color.rgb = COLOR_SECONDARY
            p.alignment = PP_ALIGN.CENTER

            current_left += widths[i]
        top += Inches(0.42)

    return slide


def add_roadmap_overview(prs):
    """実行ロードマップ概要"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "実行ロードマップ（36ヶ月）"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    phases = [
        {
            "phase": "Phase 1",
            "period": "Month 0-6",
            "title": "基盤構築・クイックウィン",
            "activities": [
                "在庫可視化、倉庫改善（1拠点）",
                "物流コスト可視化",
                "早期成果創出",
            ],
            "effect": "150-300百万円/年",
            "color": RGBColor(255, 240, 240)
        },
        {
            "phase": "Phase 2",
            "period": "Month 6-18",
            "title": "コアシステム導入・展開",
            "activities": [
                "需要予測、適正在庫基準",
                "配送最適化、VMI/IoTパイロット",
                "システム選定・要件定義",
            ],
            "effect": "500-700百万円/年",
            "color": RGBColor(255, 250, 230)
        },
        {
            "phase": "Phase 3",
            "period": "Month 18-36",
            "title": "全社展開・定着化",
            "activities": [
                "システム本格稼働（WMS/TMS）",
                "全拠点展開、VMI/IoT拡大",
                "自走体制確立",
            ],
            "effect": "870-1,240百万円/年",
            "color": RGBColor(240, 255, 240)
        }
    ]

    top = Inches(1.2)
    height = Inches(1.9)

    for phase in phases:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = phase["color"]
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{phase['phase']}：{phase['title']}　（{phase['period']}）"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(8)

        for activity in phase['activities']:
            p = tf.add_paragraph()
            p.text = f"• {activity}"
            p.font.size = Pt(13)
            p.space_before = Pt(4)

        p = tf.add_paragraph()
        p.text = f"期待効果：{phase['effect']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(10)

        top += Inches(2.05)

    return slide


def add_expected_benefits(prs):
    """期待効果まとめ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "期待効果まとめ"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 3つのカテゴリー
    categories = [
        {
            "title": "財務指標の改善",
            "items": [
                "CFマージン：-3.12% → 3.5-4.0%",
                "経常利益率：1.73% → 3.0-3.5%",
                "総利益率：7.38% → 9.0-10.0%",
                "流動比率：91.06% → 120-130%",
            ],
            "color": RGBColor(240, 255, 240)
        },
        {
            "title": "業務改善効果",
            "items": [
                "在庫削減：1,000-2,000百万円",
                "物流コスト削減：3,000-4,000百万円/年",
                "運転資金解放：5-7億円",
                "倉庫生産性向上：30-40%",
            ],
            "color": RGBColor(240, 248, 255)
        },
        {
            "title": "経営基盤の強化",
            "items": [
                "データドリブン経営の実現",
                "SCMの可視化・最適化",
                "競争力強化",
                "事業成長の基盤構築",
            ],
            "color": RGBColor(255, 250, 240)
        }
    ]

    left_start = Inches(0.5)
    width = Inches(3)
    height = Inches(5.8)

    for i, cat in enumerate(categories):
        left = left_start + i * Inches(3.15)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.2), width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = cat["color"]
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = cat["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

        for item in cat["items"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.space_before = Pt(8)

    return slide


# 協業プロジェクト計画スライド

def add_collaboration_approach(prs):
    """協業アプローチ（改善版）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "協業アプローチ：テーマ別段階的実現"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # リード文
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 240, 250)
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "現場実務に深く入り込み、テーマ別に段階的に成果を創出。\nクイックウィンで早期効果を実証し、確実に物流改革を実現します。"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # 基本方針（左）
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.85), Inches(4.3), Inches(4.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_LIGHT
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    p.text = "■ 基本方針"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    approaches = [
        "✓ テーマ別プロジェクトで段階的に成果創出",
        "✓ 現場実務に深く入り込み、実態を理解",
        "✓ クイックウィンで早期効果を実証",
        "✓ パイロット→検証→横展開",
        "✓ 貴社メンバーと協働、ノウハウ移転",
    ]

    for item in approaches:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(10)

    # プロジェクト期間（右）
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.85), Inches(4.3), Inches(4.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_LIGHT2
    shape.line.color.rgb = COLOR_SECONDARY
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    p.text = "■ プロジェクト期間：36ヶ月"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.space_after = Pt(10)

    years = [
        ("Year 1（Month 1-6）", "現状可視化、クイックウィン", "150-300百万円/年"),
        ("Year 2（Month 7-18）", "基盤構築、パイロット実施", "500-700百万円/年"),
        ("Year 3（Month 19-36）", "全社展開、自走体制確立", "870-1,240百万円/年")
    ]

    for year, desc, effect in years:
        p = tf.add_paragraph()
        p.text = year
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(10)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.space_before = Pt(3)

        p = tf.add_paragraph()
        p.text = f"効果：{effect}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(3)

    return slide


def add_10_themes_visual(prs):
    """10テーマビジュアル（2スライドに分割）"""
    # スライド1：テーマ1-5
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "10の重点テーマ（1/2）"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    themes_1 = [
        ("1", "在庫可視化", "滞留在庫特定、即時削減200-500百万円"),
        ("2", "需要予測向上", "誤差±15%→±5-8%、在庫削減300-600百万円"),
        ("3", "適正在庫基準", "回転日数7-10日短縮、在庫削減800-1,500百万円"),
        ("4", "配送最適化", "配送コスト15-20%削減、500-800百万円/年"),
        ("5", "倉庫改善", "生産性20-30%向上、人件費50-100百万円削減"),
    ]

    top = Inches(1.1)
    for num, title, effect in themes_1:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = f"テーマ{num}：{title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = f"効果：{effect}"
        p.font.size = Pt(13)
        p.space_before = Pt(3)

        top += Inches(1.25)

    # スライド2：テーマ6-10
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "10の重点テーマ（2/2）"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    themes_2 = [
        ("6", "物流コスト可視化", "コスト構造把握、改善ターゲット特定"),
        ("7", "VMIパイロット", "在庫20-30%削減、欠品率50%削減"),
        ("8", "IoT・デジタル化", "リアルタイム可視化、配送効率化"),
        ("9", "人材育成", "スキル向上、多能工化、属人化解消"),
        ("10", "KPIダッシュボード", "経営可視化、データドリブン経営実現"),
    ]

    top = Inches(1.1)
    for num, title, effect in themes_2:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = f"テーマ{num}：{title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = f"効果：{effect}"
        p.font.size = Pt(13)
        p.space_before = Pt(3)

        top += Inches(1.25)

    return slide


def add_project_timeline_visual(prs):
    """プロジェクトタイムライン（視覚的）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "プロジェクト全体タイムライン（36ヶ月）"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    phases = [
        {
            "title": "Phase 1：現状分析・クイックウィン",
            "period": "Month 0-6",
            "activities": "在庫可視化、倉庫改善、早期成果創出",
            "effect": "150-300百万円/年",
            "color": RGBColor(255, 235, 235)
        },
        {
            "title": "Phase 2：基盤構築・パイロット",
            "period": "Month 6-18",
            "activities": "需要予測、配送最適化、VMI/IoT、システム選定",
            "effect": "500-700百万円/年",
            "color": RGBColor(255, 248, 220)
        },
        {
            "title": "Phase 3：全社展開・定着化",
            "period": "Month 18-36",
            "activities": "システム導入、全拠点展開、自走体制確立",
            "effect": "870-1,240百万円/年",
            "color": RGBColor(235, 255, 235)
        }
    ]

    top = Inches(1.2)
    for phase in phases:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = phase["color"]
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{phase['title']}　（{phase['period']}）"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = phase['activities']
        p.font.size = Pt(13)
        p.space_before = Pt(4)

        p = tf.add_paragraph()
        p.text = f"累計効果：{phase['effect']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(8)

        top += Inches(1.95)

    return slide


def add_standard_process_visual(prs):
    """標準プロセス（フロー形式）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "標準活動プロセス（8ステップ）"
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    steps = [
        "1. キックオフ・計画",
        "2. 現状調査・データ収集",
        "3. 分析・課題整理",
        "4. 改善案設計",
        "5. 【意思決定】承認",
        "6. パイロット実施",
        "7. 【Go/No Go】判断",
        "8. 横展開・定着化"
    ]

    # 4x2グリッド
    positions = [
        (Inches(0.6), Inches(1.2)),
        (Inches(2.75), Inches(1.2)),
        (Inches(4.9), Inches(1.2)),
        (Inches(7.05), Inches(1.2)),
        (Inches(0.6), Inches(4.0)),
        (Inches(2.75), Inches(4.0)),
        (Inches(4.9), Inches(4.0)),
        (Inches(7.05), Inches(4.0))
    ]

    width = Inches(1.9)
    height = Inches(2.5)

    for i, (step, pos) in enumerate(zip(steps, positions)):
        # ボックス
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], width, height)
        shape.fill.solid()
        if "【" in step:
            shape.fill.fore_color.rgb = RGBColor(255, 240, 240)
            shape.line.color.rgb = COLOR_SECONDARY
            shape.line.width = Pt(3)
        else:
            shape.fill.fore_color.rgb = COLOR_BG_LIGHT
            shape.line.color.rgb = COLOR_PRIMARY
            shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.8)

        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(13)
        p.font.bold = True
        if "【" in step:
            p.font.color.rgb = COLOR_SECONDARY
        else:
            p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # 矢印（横方向）
        if i < 3 or (i >= 4 and i < 7):
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                pos[0] + width, pos[1] + height/2 - Inches(0.15),
                Inches(0.25), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_PRIMARY
            arrow.line.fill.background()

    return slide


def add_project_structure_visual(prs):
    """プロジェクト推進体制"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "プロジェクト推進体制"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 3階層
    layers = [
        {
            "title": "ステアリングコミッティ（月1回）",
            "role": "重要意思決定、進捗確認",
            "members": "社長、役員、PMO",
            "color": RGBColor(255, 230, 230),
            "height": Inches(1.3)
        },
        {
            "title": "プロジェクトマネジメントオフィス（週1回）",
            "role": "全体統括、進捗管理、課題管理",
            "members": "プロジェクトリーダー、各テーマリーダー",
            "color": RGBColor(255, 245, 220),
            "height": Inches(1.5)
        },
        {
            "title": "テーマ別ワーキンググループ（週1-2回）",
            "role": "テーマ別の詳細検討・実行（10チーム）",
            "members": "現場責任者・実務担当者、コンサルタント",
            "color": RGBColor(240, 255, 240),
            "height": Inches(1.8)
        }
    ]

    top = Inches(1.2)
    for layer in layers:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, Inches(8), layer["height"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = layer["color"]
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = layer["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = f"役割：{layer['role']}"
        p.font.size = Pt(13)
        p.space_before = Pt(4)

        p = tf.add_paragraph()
        p.text = f"メンバー：{layer['members']}"
        p.font.size = Pt(13)
        p.space_before = Pt(4)

        top += layer["height"] + Inches(0.2)

        # 矢印
        if top < Inches(6.5):
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.75), top - Inches(0.15),
                Inches(0.5), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_PRIMARY
            arrow.line.fill.background()

    return slide


def add_cumulative_effects_visual(prs):
    """累積効果の推移（視覚的）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "累積効果の推移"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    years = [
        {
            "year": "Year 1",
            "phase": "クイックウィン創出期",
            "effects": [
                "在庫削減：200-500百万円",
                "倉庫改善：30-50百万円/年",
                "配送改善：50-100百万円/年",
            ],
            "total": "150-300百万円/年",
            "cf": "CF：△50-100百万円",
            "color": RGBColor(255, 240, 240)
        },
        {
            "year": "Year 2",
            "phase": "本格展開開始期",
            "effects": [
                "需要予測・適正在庫：+400-700百万円",
                "配送最適化：+200-300百万円/年",
                "VMI/倉庫展開：+150-250百万円/年",
            ],
            "total": "500-700百万円/年",
            "cf": "CF：+50-200百万円（回収開始）",
            "color": RGBColor(255, 248, 220)
        },
        {
            "year": "Year 3",
            "phase": "フル効果達成期",
            "effects": [
                "システム稼働：+200-300百万円/年",
                "全拠点展開：+100-150百万円/年",
                "VMI/IoT拡大：+70-150百万円/年",
            ],
            "total": "870-1,240百万円/年",
            "cf": "CF：+720-1,090百万円",
            "color": RGBColor(235, 255, 235)
        }
    ]

    top = Inches(1.1)
    for year in years:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = year["color"]
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = f"{year['year']}：{year['phase']}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(6)

        for effect in year['effects']:
            p = tf.add_paragraph()
            p.text = f"• {effect}"
            p.font.size = Pt(12)
            p.space_before = Pt(2)

        p = tf.add_paragraph()
        p.text = f"累計効果：{year['total']}　{year['cf']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(8)

        top += Inches(2.05)

    return slide


def add_success_factors_visual(prs):
    """成功の5つの鍵"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "成功の5つの鍵"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    factors = [
        ("1", "経営層の強いコミットメント", "トップダウン推進、明確な目標、リソース確保"),
        ("2", "現場を巻き込んだ推進", "現場の声を反映、早期成果で実感創出"),
        ("3", "データドリブンの徹底", "事実に基づく課題把握、定量効果測定"),
        ("4", "クイックウィンの創出", "早期成果実証、組織の信頼獲得"),
        ("5", "人材育成とノウハウ移転", "貴社メンバー成長、内製化、自走体制"),
    ]

    top = Inches(1.2)
    for num, title, desc in factors:
        # 番号円
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), top + Inches(0.25), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_SECONDARY
        circle.line.fill.background()

        circle_tf = circle.text_frame
        circle_p = circle_tf.paragraphs[0]
        circle_p.text = num
        circle_p.alignment = PP_ALIGN.CENTER
        circle_p.font.size = Pt(24)
        circle_p.font.bold = True
        circle_p.font.color.rgb = RGBColor(255, 255, 255)

        # コンテンツボックス
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), top, Inches(8.1), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.08)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.space_before = Pt(2)

        top += Inches(1.15)

    return slide


def add_next_steps_visual(prs):
    """次のステップ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "次のステップ：プロジェクト開始まで"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    steps = [
        {
            "num": "Step 1",
            "title": "詳細ヒアリング・現地視察",
            "period": "2-3週間",
            "details": [
                "経営層・物流部門ヒアリング",
                "拠点視察（2-3拠点）",
                "データ確認・収集",
            ],
            "output": "詳細現状分析、優先テーマ案、詳細計画"
        },
        {
            "num": "Step 2",
            "title": "提案プレゼンテーション",
            "period": "1週間",
            "details": [
                "詳細計画のご説明",
                "期待効果の精緻化",
                "体制・役割分担の確認",
            ],
            "output": "最終提案書、契約条件合意"
        },
        {
            "num": "Step 3",
            "title": "キックオフ",
            "period": "1週間",
            "details": [
                "キックオフミーティング",
                "プロジェクト体制発足",
                "活動開始",
            ],
            "output": "プロジェクト始動"
        }
    ]

    top = Inches(1.2)
    for step in steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(9), Inches(1.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_LIGHT
        shape.line.color.rgb = COLOR_PRIMARY
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = f"{step['num']}：{step['title']}　（{step['period']}）"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(6)

        for detail in step['details']:
            p = tf.add_paragraph()
            p.text = f"• {detail}"
            p.font.size = Pt(13)
            p.space_before = Pt(3)

        p = tf.add_paragraph()
        p.text = f"成果物：{step['output']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(8)

        top += Inches(2.0)

    return slide


def add_thank_you_slide(prs):
    """Thank youスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.5), Inches(10), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 245, 250)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.text = "ご清聴ありがとうございました"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "ご質問・ご相談はお気軽にお申し付けください"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXT

    return slide


if __name__ == "__main__":
    prs = create_presentation()
    prs.save("物流ソリューション提案書_ヤマエ久野_完全版v2.pptx")
    print(f"PowerPointプレゼンテーション（改善版）を作成しました")
    print(f"総スライド数: {len(prs.slides)}")
